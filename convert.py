# -*- encoding: utf-8 -*-
import fitz
import sys
import glob
import comtypes.client
import os

from pptx import Presentation
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM


def convert_pdf_to_png(pdfFileName, pngFileName):
    doc = fitz.open(pdfFileName)

    page = doc[0]
    rotate = int(0)
    zoom_x = 3.0
    zoom_y = 3.0
    trans = fitz.Matrix(zoom_x, zoom_y).preRotate(rotate)
    pm = page.getPixmap(matrix=trans, alpha=False)
    pm.writePNG(pngFileName)


def convert_pptx_to_pdf(inputFileName, outputFileName, formatType=32):
    ppt = comtypes.client.CreateObject("Powerpoint.Application")
    ppt.Visible = 0

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = ppt.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType)
    deck.close()

    ppt.Quit()

def replace_text_slide(slide, maze_id, fn):
    pass

def replace_img_slide(slide, img, img_path):
    # Replace the picture in the shape object (img) with the image in img_path.

    imgPic = img._pic
    imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
    imgPart = slide.part.related_parts[imgRID]

    with open(img_path, 'rb') as f:
        rImgBlob = f.read()

    # replace
    imgPart._blob = rImgBlob


def convert_svg_to_maze(svgFilename, outputFileName, maze_id):
    prs = Presentation(r'C:\Users\fjort\OneDrive\Documents\mazetplt.pptx')
    tmpPrsFile = r'C:\Users\fjort\PycharmProjects\maze\temp.pptx'
    tmpPDFFile = r'C:\Users\fjort\PycharmProjects\maze\tempPdf.pdf'

    sld = prs.slides[0]
    img = sld.shapes[0]
    # fpath = r'C:\Users\fjort\Downloads\maze.svg'

    # targetDir = r"C:\Users\fjort\Downloads"
    #
    # exportFile = os.path.join(targetDir, fpath[:-3]+"png")
    # print(exportPath)

    # Convert SVG to PNG
    drawing = svg2rlg(svgFilename)
    renderPM.drawToFile(drawing, outputFileName, fmt="PNG")

    # Replace maze in template
    replace_img_slide(sld, img, outputFileName)
    replace_text_slide(sld, maze_id, '')
    prs.save(tmpPrsFile)

    # Convert ppt to pdf
    convert_pptx_to_pdf(tmpPrsFile, tmpPDFFile)

    # Convert pdf to png
    convert_pdf_to_png(tmpPDFFile, outputFileName)

    print('Convert complete')


if __name__ == '__main__':
    svg_dir = r'E:\Mazes\SVG'
    png_dir = r'E:\Mazes\PNG'

    for fn in os.listdir(svg_dir):
        # print(fn)
        # print(os.path.splitext(fn)[0])
        mazeID = os.path.splitext(fn)[0]

        svgFile = os.path.join(svg_dir, fn)
        mazeFile = os.path.join(png_dir, mazeID + '.png')
        print(mazeFile)

        convert_svg_to_maze(svgFile, mazeFile, mazeID)
        print('%s is converted.' % fn)
